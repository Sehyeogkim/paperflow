from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
A="fig5_assets/"

# palette
NAVY=RGBColor(0x21,0x40,0x6e); LBLUE=RGBColor(0xDD,0xE8,0xF5)
DGRN=RGBColor(0x2E,0x5E,0x3A); LGRN=RGBColor(0xE7,0xF0,0xE7)
PURP=RGBColor(0x4B,0x2E,0x83); LPUR=RGBColor(0xEF,0xEA,0xF7)
TAN =RGBColor(0xEC,0xE5,0xD5); GRAY=RGBColor(0x9A,0x90,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x22,0x22,0x22)
SLOT=RGBColor(0xF2,0xF0,0xEA); SLOTLN=RGBColor(0xC8,0xBF,0xAE)

prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6])
shp=s.shapes

def box(x,y,w,h,fill,line,lw=1.5,rad=0.06):
    b=shp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb=fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb=line; b.line.width=Pt(lw)
    b.shadow.inherit=False
    try: b.adjustments[0]=rad
    except: pass
    return b

def setbase(run,val):
    run._r.get_or_add_rPr().set('baseline',str(val))

def para(tf,parts,align=PP_ALIGN.LEFT,space=2,first=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0)
    for text,fmt in parts:
        r=p.add_run(); r.text=text; f=r.font
        f.size=Pt(fmt.get('sz',10)); f.bold=fmt.get('b',False)
        f.name="Arial"; f.color.rgb=fmt.get('c',INK)
        if fmt.get('sup'): setbase(r,30000)
        if fmt.get('sub'): setbase(r,-25000)
    return p

def tbox(x,y,w,h,anchor=MSO_ANCHOR.TOP):
    t=shp.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=t.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(3); tf.margin_right=Pt(3); tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    return tf

def arrow(kind,x,y,w,h,color):
    a=shp.add_shape(kind,Inches(x),Inches(y),Inches(w),Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb=color; a.line.fill.background(); a.shadow.inherit=False
    return a

def circle(x,y,d,fill,line=None):
    c=shp.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb=fill
    if line is None: c.line.fill.background()
    else: c.line.color.rgb=line; c.line.width=Pt(1.2)
    c.shadow.inherit=False; return c

def pic(path,x,y,w,h):
    if os.path.exists(path): shp.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))

def slot(x,y,w,h,label):
    box(x,y,w,h,SLOT,SLOTLN,lw=1.0,rad=0.08)
    tf=tbox(x,y,w,h,MSO_ANCHOR.MIDDLE); para(tf,[(label,{'sz':7.5,'c':GRAY})],PP_ALIGN.CENTER,first=True)

# column geometry
COLS=[0.30,4.66,9.02]; CW=4.0
BANNER=(0.25,0.70); TOPBOX=(1.12,2.78); BOTBOX=(4.45,2.78)

# ---------- BANNERS ----------
banners=[("1","Dataset Generation"),("2","Vulnerability Index Selection"),("3","Surrogate Model &\nSensitivity Analysis")]
for i,(num,title) in enumerate(banners):
    x=COLS[i]; box(x,BANNER[0],CW,BANNER[1],TAN,None,rad=0.18)
    circle(x+0.18,BANNER[0]+0.13,0.44,NAVY)
    tf=tbox(x+0.18,BANNER[0]+0.10,0.44,0.48,MSO_ANCHOR.MIDDLE); para(tf,[(num,{'sz':16,'b':True,'c':WHITE})],PP_ALIGN.CENTER,first=True)
    tf=tbox(x+0.72,BANNER[0],CW-0.8,BANNER[1],MSO_ANCHOR.MIDDLE)
    for j,ln in enumerate(title.split("\n")): para(tf,[(ln,{'sz':14,'b':True,'c':INK})],PP_ALIGN.LEFT,first=(j==0))

# ===================== STAGE 1 (blue) =====================
x=COLS[0]
box(x,TOPBOX[0],CW,TOPBOX[1],LBLUE,NAVY,1.6)
tf=tbox(x,TOPBOX[0]+0.06,CW,0.3); para(tf,[("Input Parameters",{'sz':12.5,'b':True,'c':NAVY})],PP_ALIGN.CENTER,first=True)
caps=["Hemodynamic\nPressure/Flow","Plaque\nMorphology","Material\nStress-Strain"]
imgs=["hemo.png","morph.png","stress.png"]
px=[x+0.12,x+1.42,x+2.72]
for cx,cap,im in zip(px,caps,imgs):
    tf=tbox(cx,TOPBOX[0]+0.42,1.18,0.42); 
    for j,ln in enumerate(cap.split("\n")): para(tf,[(ln,{'sz':7.8,'b':True,'c':INK})],PP_ALIGN.CENTER,first=(j==0))
    pic(A+im,cx,TOPBOX[0]+0.95,1.18,1.62)
# down arrow
arrow(MSO_SHAPE.DOWN_ARROW,x+CW/2-0.22,TOPBOX[0]+TOPBOX[1]+0.02,0.44,0.5,NAVY)
# bottom box
box(x,BOTBOX[0],CW,BOTBOX[1],LBLUE,NAVY,1.6)
tf=tbox(x+0.15,BOTBOX[0]+0.10,CW-0.3,0.3); para(tf,[("Cost-effective FSI simulation",{'sz':12,'b':True,'c':NAVY})],PP_ALIGN.LEFT,first=True)
tf=tbox(x+0.15,BOTBOX[0]+0.48,CW-0.3,0.3); para(tf,[("→ 1,000-sample dataset",{'sz':10})],first=True)
slot(x+0.18,BOTBOX[0]+0.95,1.7,1.6,"[ FSI result\nimage ]")
tf=tbox(x+2.0,BOTBOX[0]+1.0,CW-2.1,1.6)
para(tf,[("Outputs:",{'sz':10,'b':True})],first=True)
para(tf,[("• PSS (peak systolic stress)",{'sz':9.3})])
para(tf,[("• ",{'sz':9.3}),("ΔPSS (change in PSS)",{'sz':9.3})])

# ===================== STAGE 2 (green) =====================
x=COLS[1]
box(x,TOPBOX[0],CW,TOPBOX[1],LGRN,DGRN,1.6)
tf=tbox(x,TOPBOX[0]+0.10,CW,0.3); para(tf,[("6 VI candidates",{'sz':13,'b':True,'c':DGRN})],PP_ALIGN.CENTER,first=True)
rows=[("VI",[("VI = stress / strength",{'sz':11})]),
      ("σ",[("stress = {PSS, ΔPSS}",{'sz':11})]),
      ("α",[("α = {0.0, 0.5, 1.0}",{'sz':11})])]
ry=TOPBOX[0]+0.7
for sym,parts in rows:
    circle(x+0.45,ry,0.4,WHITE,DGRN)
    tf=tbox(x+0.45,ry-0.02,0.4,0.44,MSO_ANCHOR.MIDDLE); para(tf,[(sym,{'sz':11,'b':True,'c':DGRN})],PP_ALIGN.CENTER,first=True)
    tf=tbox(x+1.0,ry-0.02,CW-1.1,0.44,MSO_ANCHOR.MIDDLE); para(tf,parts,first=True)
    ry+=0.62
arrow(MSO_SHAPE.DOWN_ARROW,x+CW/2-0.22,TOPBOX[0]+TOPBOX[1]+0.02,0.44,0.5,DGRN)
box(x,BOTBOX[0],CW,BOTBOX[1],LGRN,DGRN,1.6)
tf=tbox(x+0.15,BOTBOX[0]+0.10,CW-0.3,0.3); para(tf,[("7-criterion clinical screening",{'sz':12,'b':True,'c':DGRN})],first=True)
slot(x+0.18,BOTBOX[0]+0.5,0.5,0.5,"")
tf=tbox(x+0.8,BOTBOX[0]+0.52,CW-0.9,0.4,MSO_ANCHOR.MIDDLE); para(tf,[("• Sign-agreement test",{'sz':10})],first=True)
ln=shp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(BOTBOX[0]+1.12),Inches(CW-0.4),Pt(1)); ln.fill.solid(); ln.fill.fore_color.rgb=DGRN; ln.line.fill.background(); ln.shadow.inherit=False
tf=tbox(x+0.2,BOTBOX[0]+1.2,CW-0.4,0.3); para(tf,[("Select two final indices:",{'sz':9.5,'b':True})],first=True)
def vi_tag(yy,subn,exp):
    t=box(x+0.2,yy,0.62,0.42,LGRN,DGRN,1.2,rad=0.25)
    tf=tbox(x+0.2,yy,0.62,0.42,MSO_ANCHOR.MIDDLE); para(tf,[("VI",{'sz':10.5,'b':True,'c':DGRN}),(subn,{'sz':8,'b':True,'c':DGRN,'sub':True})],PP_ALIGN.CENTER,first=True)
    tf=tbox(x+0.95,yy,CW-1.05,0.42,MSO_ANCHOR.MIDDLE)
    para(tf,[("VI",{'sz':10.5}),(subn,{'sz':7.5,'sub':True}),(" = ΔPSS / E",{'sz':10.5}),("FC",{'sz':7.5,'sub':True}),(exp,{'sz':8.5,'sup':True})],first=True)
vi_tag(BOTBOX[0]+1.6,"1","0.5"); vi_tag(BOTBOX[0]+2.12,"2","1.0")

# ===================== STAGE 3 (purple) =====================
x=COLS[2]
box(x,TOPBOX[0],CW,TOPBOX[1],LPUR,PURP,1.6)
tf=tbox(x,TOPBOX[0]+0.08,CW,0.3); para(tf,[("GPR surrogate model",{'sz':12.5,'b':True,'c':PURP})],PP_ALIGN.CENTER,first=True)
tf=tbox(x,TOPBOX[0]+0.40,CW,0.3); para(tf,[("Map input parameters → VI",{'sz':9.5})],PP_ALIGN.CENTER,first=True)
tf=tbox(x+0.12,TOPBOX[0]+0.85,1.6,1.7)
for j,(lab,col) in enumerate([("Hemodynamic",NAVY),("Morphological",RGBColor(0xc0,0x39,0x2b)),("Material",DGRN)]):
    para(tf,[("● ",{'sz':10,'c':col}),(lab,{'sz':8.8})],first=(j==0),space=6)
pic(A+"gpr.png",x+1.75,TOPBOX[0]+0.8,2.1,1.6)
arrow(MSO_SHAPE.DOWN_ARROW,x+CW/2-0.22,TOPBOX[0]+TOPBOX[1]+0.02,0.44,0.5,PURP)
box(x,BOTBOX[0],CW,BOTBOX[1],LPUR,PURP,1.6)
tf=tbox(x+0.15,BOTBOX[0]+0.10,CW-0.3,0.3); para(tf,[("Sobol sensitivity analysis",{'sz':12,'b':True,'c':PURP})],first=True)
pic(A+"sobol.png",x+0.7,BOTBOX[0]+0.5,2.6,1.55)
cb=box(x+0.25,BOTBOX[0]+2.2,CW-0.5,0.45,RGBColor(0xE3,0xDC,0xF0),PURP,1.0,rad=0.3)
tf=tbox(x+0.25,BOTBOX[0]+2.2,CW-0.5,0.45,MSO_ANCHOR.MIDDLE)
para(tf,[("★  ",{'sz':10,'c':PURP}),("Material > Hemodynamic > Morphological",{'sz':9.5,'b':True,'c':PURP})],PP_ALIGN.CENTER,first=True)

# ---------- inter-stage horizontal arrows ----------
ay=BOTBOX[0]+1.0
arrow(MSO_SHAPE.RIGHT_ARROW,COLS[0]+CW+0.02,ay,0.30,0.5,NAVY)
arrow(MSO_SHAPE.RIGHT_ARROW,COLS[1]+CW+0.02,ay,0.30,0.5,DGRN)

out="fig5_editable.pptx"; prs.save(out); print("saved",out, os.path.getsize(out),"bytes")
